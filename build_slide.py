"""
Recreates the Sunline Tech Platform - APIBase architecture slide as an editable .pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0d, 0x1c, 0x30)
DARK_BLUE   = RGBColor(0x14, 0x2b, 0x46)
MED_BLUE    = RGBColor(0x1c, 0x42, 0x68)
STEEL       = RGBColor(0x2c, 0x60, 0x90)
LIGHT_STEEL = RGBColor(0x9c, 0xc0, 0xd8)
TPC_BG      = RGBColor(0xc2, 0xdb, 0xea)
BOX_BG      = RGBColor(0xec, 0xf5, 0xfb)
WHITE       = RGBColor(0xff, 0xff, 0xff)
GRAY_LINE   = RGBColor(0xb0, 0xc8, 0xd8)
DARK_TEXT   = RGBColor(0x1a, 0x1a, 0x1a)
ORANGE      = RGBColor(0xf0, 0x90, 0x00)
TEAL_HDR    = RGBColor(0x00, 0x7a, 0x8a)

# ── Helpers ────────────────────────────────────────────────────────────────────
def R(l, t, w, h, fill, line=None, lw=0.5):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else:    s.line.fill.background()
    return s

def RT(l, t, w, h, text, fill, line=None, lw=0.5,
       fs=7, bold=False, fc=WHITE, align=PP_ALIGN.CENTER, wrap=True, anchor=None):
    s = R(l, t, w, h, fill, line, lw)
    tf = s.text_frame; tf.word_wrap = wrap
    if anchor: tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(fs); run.font.bold = bold; run.font.color.rgb = fc
    return s

def TB(l, t, w, h, text, fs=7, bold=False, fc=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(fs); run.font.bold = bold; run.font.color.rgb = fc
    return tb

def small_box(l, t, w, h, text, fs=5.5):
    RT(l, t, w, h, text, BOX_BG, line=GRAY_LINE, lw=0.3,
       fs=fs, fc=DARK_TEXT, align=PP_ALIGN.CENTER, wrap=True)

# ── 1. Background ──────────────────────────────────────────────────────────────
R(0, 0, 13.33, 7.5, NAVY)

# ── 2. Partner Ecosystem banner (full-width top row) ───────────────────────────
R(0, 0, 13.33, 0.44, DARK_BLUE)
RT(0.02, 0.03, 1.1, 0.38, "Partner\nEcosystem", MED_BLUE,
   fs=7, bold=True, fc=WHITE, align=PP_ALIGN.CENTER)

partners = ["Ecommerce","Education","Hospitality","Healthcare","Transportation","Recreation","ETL"]
for i, p in enumerate(partners):
    RT(1.18 + i * 1.25, 0.05, 1.18, 0.34, p, STEEL,
       line=LIGHT_STEEL, lw=0.5, fs=6.5, fc=WHITE)
# extra block on right
RT(10.0, 0.05, 1.3, 0.34, "[more icons]", STEEL, line=LIGHT_STEEL, lw=0.5, fs=6, fc=WHITE)

# ── 3. Main center frame + Sunline Tech Platform header ───────────────────────
# outer border for the whole platform
R(1.14, 0.44, 10.36, 6.32, MED_BLUE, line=STEEL, lw=1)
RT(1.14, 0.44, 10.36, 0.30, "Sunline Tech Platform  –  APIBase",
   MED_BLUE, fs=8, bold=True, fc=WHITE)

# ── 4. Transaction Processing Center ──────────────────────────────────────────
RT(1.14, 0.74, 10.36, 0.25, "Transaction Processing Center",
   TEAL_HDR, fs=7.5, bold=True, fc=WHITE)
R(1.14, 0.99, 10.36, 4.12, TPC_BG, line=STEEL, lw=0.4)

# "Joint Launching Available" – vertical label on left edge of TPC
tb_jl = TB(1.15, 1.35, 0.28, 2.5,
           "Joint Launching Available", fs=6.5, bold=True, fc=WHITE)
tb_jl.rotation = 270

# ─── 4a. Customer Center ───────────────────────────────────────────────────────
CX, CY, CW, CH = 1.55, 1.02, 2.12, 3.02
RT(CX, CY, CW, 0.22, "Customer Center", MED_BLUE, fs=7, bold=True, fc=WHITE)
R(CX, CY + 0.22, CW, CH - 0.22, WHITE, line=GRAY_LINE, lw=0.4)

RT(CX + 0.02, CY + 0.25, 1.03, 0.20, "Individual", STEEL, fs=6, fc=WHITE)
RT(CX + 1.07, CY + 0.25, 1.02, 0.20, "Corporate",  STEEL, fs=6, fc=WHITE)

cc_items = [
    "Requirement Management","Credit Limit",
    "Collateral","Authorization",
    "Document","Questionnaire",
    "KYC Verification","Black List",
    "Customer Group","Customer Product",
]
for idx, item in enumerate(cc_items):
    r_, c_ = divmod(idx, 2)
    small_box(CX + 0.03 + c_ * 1.03, CY + 0.48 + r_ * 0.26, 1.00, 0.23, item)

# ─── 4b. Loan Center ───────────────────────────────────────────────────────────
LX, LY, LW, LH = 3.70, 1.02, 1.90, 3.02
RT(LX, LY, LW, 0.22, "Loan Center", MED_BLUE, fs=7, bold=True, fc=WHITE)
R(LX, LY + 0.22, LW, LH - 0.22, WHITE, line=GRAY_LINE, lw=0.4)

RT(LX + 0.02, LY + 0.25, 0.88, 0.20, "Primary Process",  STEEL, fs=5.5, fc=WHITE, wrap=True)
RT(LX + 0.93, LY + 0.25, 0.93, 0.20, "Secondary",         STEEL, fs=5.5, fc=WHITE)

loan_items = ["Guidance","Application","Credit Process",
              "Loan Booking","Disbursement","Settlement","Penalty"]
for idx, item in enumerate(loan_items):
    small_box(LX + 0.03, LY + 0.48 + idx * 0.26, LW - 0.06, 0.23, item)

# Factories sub-row at bottom of loan
factories = ["Pricing Factory","Product Factory","Exchange Rate","Debit Card"]
for idx, item in enumerate(factories):
    c_ = idx % 2; r_ = idx // 2
    small_box(LX + 0.03 + c_ * 0.91, LY + 2.38 + r_ * 0.26, 0.88, 0.23, item)

# ─── 4c. Account ───────────────────────────────────────────────────────────────
AX, AY, AW, AH = 5.63, 1.02, 1.10, 1.55
RT(AX, AY, AW, 0.22, "Account", MED_BLUE, fs=7, bold=True, fc=WHITE)
R(AX, AY + 0.22, AW, AH - 0.22, WHITE, line=GRAY_LINE, lw=0.4)
for idx, item in enumerate(["Savings Product","Demand Deposit","Time Deposit","Overdraft"]):
    small_box(AX + 0.03, AY + 0.25 + idx * 0.27, AW - 0.06, 0.24, item)

# ─── 4d. Treasury / Trade Finance Center ──────────────────────────────────────
TX, TY, TW, TH = 6.76, 1.02, 1.92, 3.02
RT(TX, TY, TW, 0.22, "Treasury / Trade Finance", MED_BLUE, fs=6.5, bold=True, fc=WHITE, wrap=True)
R(TX, TY + 0.22, TW, TH - 0.22, WHITE, line=GRAY_LINE, lw=0.4)

RT(TX + 0.02, TY + 0.25, 0.90, 0.20, "Treasury",   STEEL, fs=6, fc=WHITE)
RT(TX + 0.95, TY + 0.25, 0.93, 0.20, "Trade Fin.",  STEEL, fs=6, fc=WHITE)

tf_items = ["FX Rate","Money Market","Bond / Securities",
            "LC / LG Issuance","Import / Export","Guarantee","Liability Mgmt"]
for idx, item in enumerate(tf_items):
    small_box(TX + 0.03, TY + 0.48 + idx * 0.26, TW - 0.06, 0.23, item)

# ─── 4e. Remittance Center ─────────────────────────────────────────────────────
RX, RY, RW, RH = 8.71, 1.02, 1.20, 2.20
RT(RX, RY, RW, 0.22, "Remittance\nCenter", MED_BLUE, fs=6.5, bold=True, fc=WHITE, wrap=True)
R(RX, RY + 0.22, RW, RH - 0.22, WHITE, line=GRAY_LINE, lw=0.4)
for idx, item in enumerate(["Domestic Transfer","Cross Border","SWIFT","ACH / RTGS","Standing Order"]):
    small_box(RX + 0.03, RY + 0.25 + idx * 0.28, RW - 0.06, 0.24, item)

# ─── 4f. Card Center ───────────────────────────────────────────────────────────
KX, KY, KW, KH = 9.95, 1.02, 1.48, 3.02
RT(KX, KY, KW, 0.22, "Card Center", MED_BLUE, fs=7, bold=True, fc=WHITE)
R(KX, KY + 0.22, KW, KH - 0.22, WHITE, line=GRAY_LINE, lw=0.4)

card_items_2col = ["OA Center","3E API","Gateway","Partner","FX Manager","Card Issuing"]
for idx, item in enumerate(card_items_2col):
    c_ = idx % 2; r_ = idx // 2
    small_box(KX + 0.03 + c_ * 0.71, KY + 0.25 + r_ * 0.27, 0.68, 0.24, item)
small_box(KX + 0.03, KY + 1.09, KW - 0.06, 0.24, "Card Acquiring")
small_box(KX + 0.03, KY + 1.36, KW - 0.06, 0.24, "Loyalty Points")

# ─── 4g. Marketing Factory (under Account, mid section) ───────────────────────
MX, MY, MW, MH = 5.63, 2.60, 1.10, 1.48
RT(MX, MY, MW, 0.22, "Marketing\nFactory", MED_BLUE, fs=6, bold=True, fc=WHITE, wrap=True)
R(MX, MY + 0.22, MW, MH - 0.22, WHITE, line=GRAY_LINE, lw=0.4)
for idx, item in enumerate(["Campaign Mgmt","Loyalty Program","Offer Engine","Analytics"]):
    small_box(MX + 0.03, MY + 0.25 + idx * 0.27, MW - 0.06, 0.24, item)

# ─── 4h. Operation Center (bottom-right of TPC) ────────────────────────────────
OX, OY, OW, OH = 8.71, 3.25, 2.72, 1.82
RT(OX, OY, OW, 0.22, "Operation Center", MED_BLUE, fs=7, bold=True, fc=WHITE)
R(OX, OY + 0.22, OW, OH - 0.22, WHITE, line=GRAY_LINE, lw=0.4)
oc_items = ["End of Day","Batch Processing","Reconciliation","GL Interface","Authority","Contract"]
for idx, item in enumerate(oc_items):
    c_ = idx % 2; r_ = idx // 2
    small_box(OX + 0.03 + c_ * 1.33, OY + 0.25 + r_ * 0.27, 1.30, 0.24, item)

# ─── 4i. Sundry / Equity / Profiling ──────────────────────────────────────────
SX, SY, SW, SH = 6.76, 3.25, 1.92, 1.82
RT(SX, SY, SW, 0.22, "Sundry Debit / Credit", MED_BLUE, fs=6.5, bold=True, fc=WHITE)
R(SX, SY + 0.22, SW, SH - 0.22, WHITE, line=GRAY_LINE, lw=0.4)
for idx, item in enumerate(["Equity Management","Profiling","Enquiry","Reporting","GL Mapping","Tax Calculation"]):
    c_ = idx % 2; r_ = idx // 2
    small_box(SX + 0.03 + c_ * 0.93, SY + 0.25 + r_ * 0.27, 0.90, 0.24, item)

# ── 5. Account Processing Center ──────────────────────────────────────────────
APCY = 5.14
R(1.14, APCY, 10.36, 0.22, MED_BLUE)
RT(1.14, APCY, 10.36, 0.22, "Account Processing Center",
   MED_BLUE, fs=7, bold=True, fc=WHITE)
R(1.14, APCY + 0.22, 10.36, 0.78, DARK_BLUE, line=STEEL, lw=0.4)

apc_items = [
    "ETL Base",
    "Protocol\n(Parsing Protocol)",
    "Message Format\nManagement Protocol",
    "Open API",
    "Routing\nManagement",
    "API Tenant",
    "AI Center",
]
for idx, item in enumerate(apc_items):
    small_box(1.19 + idx * 1.44, APCY + 0.25, 1.37, 0.66, item, fs=6)

# ── 6. Left Access Point panel ────────────────────────────────────────────────
R(0, 0.44, 1.14, 6.45, DARK_BLUE, line=STEEL, lw=0.5)
RT(0.02, 0.46, 1.10, 0.26, "Access Point", MED_BLUE, fs=7, bold=True, fc=WHITE)

access_items = [
    "Customer\nCenter", "Branch\nChannel", "ATM",
    "Cash\nManagement", "IVR",
    "AFM / Z-Box /\nSmart Box",
    "Internet\nBanking", "Fuel\nLimiting", "APIs",
]
for idx, item in enumerate(access_items):
    RT(0.05, 0.76 + idx * 0.60, 1.04, 0.52,
       item, MED_BLUE, line=LIGHT_STEEL, lw=0.4,
       fs=6, fc=WHITE, align=PP_ALIGN.CENTER, wrap=True)

RT(0.05, 6.10, 0.50, 0.28, "Sunline", ORANGE, fs=7, bold=True, fc=WHITE)
RT(0.57, 6.10, 0.54, 0.28, "Partner", MED_BLUE, line=LIGHT_STEEL, lw=0.4, fs=7, fc=WHITE)

# ── 7. Right Extended Service panel ───────────────────────────────────────────
R(11.50, 0.44, 1.83, 6.45, DARK_BLUE, line=STEEL, lw=0.5)
RT(11.52, 0.46, 1.58, 0.26, "Extended Service", MED_BLUE, fs=7, bold=True, fc=WHITE)

ext_items = [
    "UAM","Origination","Collection",
    "Interactive\nDocuments","Trouble Desk",
    "Digital Lending","Smart Analytics",
    "Contract","CRM","Law, Limit",
]
for idx, item in enumerate(ext_items):
    RT(11.54, 0.76 + idx * 0.55, 1.10, 0.48,
       item, MED_BLUE, line=LIGHT_STEEL, lw=0.4,
       fs=6, fc=WHITE, align=PP_ALIGN.CENTER, wrap=True)

# "Enterprise Integration Layer" – vertical text on far right edge
tb_eil = TB(13.02, 1.20, 0.28, 4.0,
            "Enterprise Integration Layer",
            fs=6.5, bold=True, fc=WHITE)
tb_eil.rotation = 270

# ── 8. Sunline Data Platform (bottom full-width row) ──────────────────────────
SDPY = 6.90
R(0, SDPY, 13.33, 0.60, RGBColor(0x08, 0x14, 0x22))
RT(0.05, SDPY + 0.02, 1.38, 0.54, "Sunline\nData Platform",
   ORANGE, fs=7, bold=True, fc=WHITE, align=PP_ALIGN.CENTER)

sdp_items = [
    "Enterprise\nData Warehouse",
    "Data Analytics",
    "Enterprise Report",
    "Precision\nMarketing",
    "Distributed\nArchitecture",
]
for idx, item in enumerate(sdp_items):
    RT(1.48 + idx * 2.15, SDPY + 0.04, 2.0, 0.52,
       item, DARK_BLUE, line=STEEL, lw=0.4,
       fs=7, fc=WHITE, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
out = "/home/bibd1611/codes/HRMSv2/sunline_architecture.pptx"
prs.save(out)
print(f"Saved → {out}")
